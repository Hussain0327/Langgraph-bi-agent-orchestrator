from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path
from typing import Optional, List
from datetime import datetime
import io
from src.schemas import AgentOutput, Recommendation
from .chart_generator import ChartGenerator

class ValtricTheme:
    PRIMARY = '#2C3E50'
    SECONDARY = '#3498DB'
    ACCENT = '#E74C3C'
    SUCCESS = '#2ECC71'
    WARNING = '#F39C12'
    TEXT = '#2C3E50'
    BACKGROUND = '#FFFFFF'

    @staticmethod
    def hex_to_rgb(hex_color: str) -> RGBColor:
        hex_color = hex_color.lstrip('#')
        return RGBColor(*[int(hex_color[i:i + 2], 16) for i in (0, 2, 4)])
    TITLE_FONT = 'Calibri'
    BODY_FONT = 'Calibri'
    TITLE_SIZE = Pt(32)
    HEADING_SIZE = Pt(24)
    BODY_SIZE = Pt(14)
    CAPTION_SIZE = Pt(10)

class PowerPointGenerator:

    def __init__(self, theme: ValtricTheme=None):
        self.theme = theme or ValtricTheme()
        self.chart_gen = ChartGenerator()

    def generate(self, agent_output: AgentOutput, output_path: Optional[str]=None, template: str='executive_summary') -> str:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        if template == 'executive_summary':
            self._build_executive_summary(prs, agent_output)
        else:
            raise ValueError(f'Unknown template: {template}')
        if not output_path:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_path = f'presentation_{agent_output.agent}_{timestamp}.pptx'
        prs.save(output_path)
        print(f'✓ PowerPoint saved: {output_path}')
        return output_path

    def _build_executive_summary(self, prs: Presentation, output: AgentOutput):
        self._add_title_slide(prs, output)
        self._add_executive_summary_slide(prs, output)
        self._add_context_slide(prs, output)
        self._add_key_findings_slide(prs, output)
        if output.findings.metrics:
            self._add_metrics_slides(prs, output)
        if output.findings.risks:
            self._add_risks_slide(prs, output)
        if output.findings.recommendations:
            self._add_recommendations_slide(prs, output)
        self._add_next_steps_slide(prs, output)
        if output.research_citations:
            self._add_appendix_slide(prs, output)

    def _add_title_slide(self, prs: Presentation, output: AgentOutput):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme.hex_to_rgb(self.theme.PRIMARY)
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = 'Business Intelligence Analysis'
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(7), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = output.query[:100] + ('...' if len(output.query) > 100 else '')
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = f"Prepared by ValtricAI | {output.timestamp.strftime('%B %d, %Y')}"
        footer_para = footer_frame.paragraphs[0]
        footer_para.alignment = PP_ALIGN.CENTER
        footer_para.font.size = Pt(12)
        footer_para.font.color.rgb = RGBColor(200, 200, 200)

    def _add_executive_summary_slide(self, prs: Presentation, output: AgentOutput):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = 'Executive Summary'
        self._style_title(title)
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = output.findings.executive_summary
        p.font.size = Pt(14)
        p.space_after = Pt(12)
        if output.findings.metrics and len(output.findings.metrics) > 0:
            p = text_frame.add_paragraph()
            p.text = '\nKey Metrics:'
            p.font.bold = True
            p.font.size = Pt(16)
            p.space_before = Pt(20)
            for i, (name, metric) in enumerate(list(output.findings.metrics.items())[:3]):
                p = text_frame.add_paragraph()
                p.text = f"• {name.replace('_', ' ').title()}: {metric.value} {metric.unit}"
                p.font.size = Pt(14)
                p.level = 1

    def _add_context_slide(self, prs: Presentation, output: AgentOutput):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = 'Context'
        self._style_title(title)
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = 'Question'
        p.font.bold = True
        p.font.size = Pt(16)
        p = text_frame.add_paragraph()
        p.text = output.query
        p.font.size = Pt(14)
        p.space_after = Pt(20)
        p = text_frame.add_paragraph()
        p.text = 'Why This Matters'
        p.font.bold = True
        p.font.size = Pt(16)
        p.space_before = Pt(20)
        p = text_frame.add_paragraph()
        p.text = f"This analysis provides {output.agent.replace('_', ' ')} insights to inform strategic business decisions."
        p.font.size = Pt(14)
        p = text_frame.add_paragraph()
        p.text = '\nAnalysis Details'
        p.font.bold = True
        p.font.size = Pt(16)
        p.space_before = Pt(20)
        details = [f"Agent: {output.agent.replace('_', ' ').title()}", f'Model: {output.metadata.model}', f'Confidence: {output.metadata.confidence.title()}', f"Generated: {output.timestamp.strftime('%B %d, %Y at %I:%M %p')}"]
        for detail in details:
            p = text_frame.add_paragraph()
            p.text = f'• {detail}'
            p.font.size = Pt(12)
            p.level = 1

    def _add_key_findings_slide(self, prs: Presentation, output: AgentOutput):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = 'Key Findings'
        self._style_title(title)
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        for i, finding in enumerate(output.findings.key_findings, 1):
            p = text_frame.paragraphs[0] if i == 1 else text_frame.add_paragraph()
            p.text = finding
            p.font.size = Pt(14)
            p.space_after = Pt(12)
            p.level = 0
            p.text = f'• {finding}'

    def _add_metrics_slides(self, prs: Presentation, output: AgentOutput):
        metrics_dict = {name.replace('_', ' ').title(): metric.value for name, metric in output.findings.metrics.items()}
        metric_items = list(metrics_dict.items())
        chunk_size = 6
        for chunk_idx in range(0, len(metric_items), chunk_size):
            chunk = dict(metric_items[chunk_idx:chunk_idx + chunk_size])
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = f'Key Metrics' + (f' (Part {chunk_idx // chunk_size + 1})' if len(metric_items) > chunk_size else '')
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = self.theme.hex_to_rgb(self.theme.PRIMARY)
            chart_bytes = self.chart_gen.generate_metric_comparison(metrics=chunk, title='', output_path=None)
            metrics_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
            metrics_frame = metrics_box.text_frame
            for name, value in chunk.items():
                p = metrics_frame.paragraphs[0] if name == list(chunk.keys())[0] else metrics_frame.add_paragraph()
                p.text = f'{name}: {value:,}' if isinstance(value, (int, float)) else f'{name}: {value}'
                p.font.size = Pt(18)
                p.space_after = Pt(15)

    def _add_risks_slide(self, prs: Presentation, output: AgentOutput):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = 'Risks & Considerations'
        self._style_title(title)
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        for i, risk in enumerate(output.findings.risks, 1):
            p = text_frame.paragraphs[0] if i == 1 else text_frame.add_paragraph()
            p.text = f'• {risk}'
            p.font.size = Pt(14)
            p.space_after = Pt(12)

    def _add_recommendations_slide(self, prs: Presentation, output: AgentOutput):
        for rec in output.findings.recommendations:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            priority_emoji = '🔴' if rec.priority == 'high' else '🟡' if rec.priority == 'medium' else '🟢'
            title = slide.shapes.title
            title.text = f'{priority_emoji} {rec.title}'
            self._style_title(title)
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = 'Expected Impact'
            p.font.bold = True
            p.font.size = Pt(16)
            p = text_frame.add_paragraph()
            p.text = rec.impact
            p.font.size = Pt(14)
            p.space_after = Pt(20)
            p = text_frame.add_paragraph()
            p.text = 'Rationale'
            p.font.bold = True
            p.font.size = Pt(16)
            p.space_before = Pt(15)
            p = text_frame.add_paragraph()
            p.text = rec.rationale
            p.font.size = Pt(14)
            p.space_after = Pt(20)
            p = text_frame.add_paragraph()
            p.text = 'Action Items'
            p.font.bold = True
            p.font.size = Pt(16)
            p.space_before = Pt(15)
            for action in rec.action_items:
                p = text_frame.add_paragraph()
                p.text = f'• {action}'
                p.font.size = Pt(13)
                p.level = 1

    def _add_next_steps_slide(self, prs: Presentation, output: AgentOutput):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = 'Next Steps'
        self._style_title(title)
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        if output.findings.recommendations:
            p = text_frame.paragraphs[0]
            p.text = 'Immediate Actions (High Priority)'
            p.font.bold = True
            p.font.size = Pt(16)
            for rec in output.findings.recommendations:
                if rec.priority == 'high':
                    p = text_frame.add_paragraph()
                    p.text = f'• {rec.title}'
                    p.font.size = Pt(14)
                    p.space_after = Pt(8)
            p = text_frame.add_paragraph()
            p.text = '\n30-Day Actions (Medium Priority)'
            p.font.bold = True
            p.font.size = Pt(16)
            p.space_before = Pt(20)
            for rec in output.findings.recommendations:
                if rec.priority == 'medium':
                    p = text_frame.add_paragraph()
                    p.text = f'• {rec.title}'
                    p.font.size = Pt(14)
                    p.space_after = Pt(8)

    def _add_appendix_slide(self, prs: Presentation, output: AgentOutput):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = 'References'
        self._style_title(title)
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        for i, citation in enumerate(output.research_citations, 1):
            p = text_frame.paragraphs[0] if i == 1 else text_frame.add_paragraph()
            authors_str = citation.authors[0] + (' et al.' if len(citation.authors) > 1 else '')
            citation_text = f'{authors_str} ({citation.year}). {citation.title}.'
            if citation.url:
                citation_text += f' {citation.url}'
            p.text = citation_text
            p.font.size = Pt(11)
            p.space_after = Pt(10)

    def _style_title(self, title_shape):
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.theme.hex_to_rgb(self.theme.PRIMARY)